"""Scaffolding script - creates all source files for the monthly report project."""
import pathlib

BASE = pathlib.Path(r'c:\Users\kseaman\Downloads\Cursor\monthly-report')

files = {}

# ─── src/data/__init__.py ───
files['src/data/__init__.py'] = '''"""Data source connectors for BigQuery and SQL Server."""

from .bigquery_source import BigQuerySource
from .sqlserver_source import SQLServerSource

__all__ = ["BigQuerySource", "SQLServerSource"]
'''

# ─── src/data/bigquery_source.py ───
files['src/data/bigquery_source.py'] = '''"""BigQuery connector for spend and other marketing data."""

import logging
from typing import Optional

import pandas as pd
from google.cloud import bigquery

logger = logging.getLogger(__name__)


class BigQuerySource:
    """Connects to BigQuery and runs parameterized queries."""

    def __init__(self, project_id: str, credentials_path: Optional[str] = None):
        self.project_id = project_id
        if credentials_path:
            self.client = bigquery.Client.from_service_account_json(
                credentials_path, project=project_id
            )
        else:
            self.client = bigquery.Client(project=project_id)

    def query(self, sql: str, params: Optional[dict] = None) -> pd.DataFrame:
        """Run a query and return results as a DataFrame."""
        job_config = bigquery.QueryJobConfig()

        if params:
            query_params = []
            for key, value in params.items():
                if isinstance(value, str):
                    query_params.append(
                        bigquery.ScalarQueryParameter(key, "STRING", value)
                    )
                elif isinstance(value, int):
                    query_params.append(
                        bigquery.ScalarQueryParameter(key, "INT64", value)
                    )
                elif isinstance(value, float):
                    query_params.append(
                        bigquery.ScalarQueryParameter(key, "FLOAT64", value)
                    )
            job_config.query_parameters = query_params

        logger.info("Executing BigQuery: %s", sql[:100])
        result = self.client.query(sql, job_config=job_config).to_dataframe()
        logger.info("Returned %d rows", len(result))
        return result

    def get_channel_spend(
        self, report_month: str, lookback_months: int = 12
    ) -> pd.DataFrame:
        """Pull channel spend data for the report window.

        Args:
            report_month: Target month as \'YYYY-MM\'
            lookback_months: How many months of history to include
        """
        sql = f"""
        SELECT
            FORMAT_DATE(\'%Y%m\', month_date) AS month_id,
            channel,
            SUM(spend) AS spend
        FROM `{self.project_id}.marketing.channel_spend`
        WHERE month_date >= DATE_SUB(PARSE_DATE(\'%Y-%m\', @report_month), INTERVAL {lookback_months} MONTH)
          AND month_date <= PARSE_DATE(\'%Y-%m\', @report_month)
        GROUP BY 1, 2
        ORDER BY 1, 2
        """
        return self.query(sql, params={"report_month": report_month})
'''

# ─── src/data/sqlserver_source.py ───
files['src/data/sqlserver_source.py'] = '''"""SQL Server / SSAS connector for lead funnel and forecast data."""

import logging
from typing import Optional

import pandas as pd
import pyodbc

logger = logging.getLogger(__name__)


class SQLServerSource:
    """Connects to SQL Server (relational or SSAS via ADOMD/OLEDB)."""

    def __init__(
        self,
        server: str,
        database: str,
        driver: str = "ODBC Driver 17 for SQL Server",
        trusted_connection: bool = True,
        username: Optional[str] = None,
        password: Optional[str] = None,
    ):
        self.server = server
        self.database = database
        self.driver = driver
        self.trusted_connection = trusted_connection
        self.username = username
        self.password = password

    def _get_connection_string(self) -> str:
        parts = [
            f"DRIVER={{{self.driver}}}",
            f"SERVER={self.server}",
            f"DATABASE={self.database}",
        ]
        if self.trusted_connection:
            parts.append("Trusted_Connection=yes")
        else:
            parts.append(f"UID={self.username}")
            parts.append(f"PWD={self.password}")
        return ";".join(parts)

    def query(self, sql: str, params: Optional[dict] = None) -> pd.DataFrame:
        """Run a SQL query and return results as a DataFrame."""
        conn_str = self._get_connection_string()
        logger.info("Connecting to %s/%s", self.server, self.database)

        with pyodbc.connect(conn_str) as conn:
            if params:
                # Simple string substitution for named params
                for key, value in params.items():
                    placeholder = f"@{key}"
                    if isinstance(value, str):
                        sql = sql.replace(placeholder, f"\\'{value}\\'")
                    else:
                        sql = sql.replace(placeholder, str(value))

            logger.info("Executing: %s", sql[:100])
            df = pd.read_sql(sql, conn)
            logger.info("Returned %d rows", len(df))
            return df

    def get_funnel_by_channel_month(self, start_month: str) -> pd.DataFrame:
        """Pull lead funnel metrics from the cube\\\'s underlying tables.

        Returns DataFrame with columns:
            year_for_month, month_name, channel, is_lead, is_application,
            is_decision_14_days, is_decision, is_new_enrollment
        """
        sql = """
        SELECT
            d.year_for_month,
            d.month_long_name,
            d.month_number_in_year,
            l.initial_marketing_segment_rollup AS channel,
            COUNT(CASE WHEN f.is_lead = 1 THEN 1 END) AS is_lead,
            COUNT(CASE WHEN f.is_application = 1 THEN 1 END) AS is_application,
            COUNT(CASE WHEN f.is_application_review = 1 THEN 1 END) AS is_application_review,
            COUNT(CASE WHEN f.is_decision_14_days = 1 THEN 1 END) AS is_decision_14_days,
            COUNT(CASE WHEN f.is_decision = 1 THEN 1 END) AS is_decision,
            COUNT(CASE WHEN f.is_new_enrollment = 1 THEN 1 END) AS is_new_enrollment
        FROM fact_lead f
        JOIN dim_date_inquiry d ON f.inquiry_date_key = d.date_key
        JOIN dim_lead l ON f.lead_key = l.lead_key
        WHERE d.year_for_month >= @start_month
        GROUP BY
            d.year_for_month, d.month_long_name, d.month_number_in_year,
            l.initial_marketing_segment_rollup
        ORDER BY d.year_for_month, d.month_number_in_year
        """
        return self.query(sql, params={"start_month": start_month})

    def get_forecast(self, month_year_id: str) -> pd.DataFrame:
        """Pull forecast data for a given month."""
        sql = """
        SELECT
            marketing_segment_rollup,
            mars_segment_new,
            day_date,
            forecasted_leads,
            MonthYearID,
            week_starting_sun,
            MTD_FLAG
        FROM forecast_table
        WHERE MonthYearID = @month_year_id
        """
        return self.query(sql, params={"month_year_id": month_year_id})

    def get_fraud_data(self, start_month: str) -> pd.DataFrame:
        """Pull fraud lead/decision counts by channel and month."""
        sql = """
        SELECT
            inquiry_month,
            initial_marketing_segment_rollup,
            Leads AS fraud_leads,
            Decisions AS fraud_decisions
        FROM fraud_view
        WHERE inquiry_month >= @start_month
        ORDER BY inquiry_month
        """
        return self.query(sql, params={"start_month": start_month})
'''

# ─── src/transforms/__init__.py ───
files['src/transforms/__init__.py'] = '''"""Metric transformations and scorecard calculations."""

from .metrics import compute_channel_metrics, compute_scorecard
from .comparisons import compute_deltas, compute_ytd

__all__ = ["compute_channel_metrics", "compute_scorecard", "compute_deltas", "compute_ytd"]
'''

# ─── src/transforms/metrics.py ───
files['src/transforms/metrics.py'] = """\"\"\"Core metric computations matching the Excel formulas.\"\"\"

import pandas as pd
import numpy as np


def compute_channel_metrics(spend: pd.DataFrame, funnel: pd.DataFrame) -> pd.DataFrame:
    \"\"\"Join spend and funnel data, compute derived metrics.

    Returns a DataFrame with one row per channel × month containing:
        spend, leads, applications, decisions_14d, enrollments,
        cpl, cpd, cpne, lead_to_decision_rate, lead_to_enroll_rate
    \"\"\"
    merged = pd.merge(
        spend,
        funnel,
        on=["month_id", "channel"],
        how="outer",
    )

    merged["cpl"] = np.where(
        merged["is_lead"] > 0,
        merged["spend"] / merged["is_lead"],
        np.nan,
    )
    merged["cpd"] = np.where(
        merged["is_decision_14_days"] > 0,
        merged["spend"] / merged["is_decision_14_days"],
        np.nan,
    )
    merged["cpne"] = np.where(
        merged["is_new_enrollment"] > 0,
        merged["spend"] / merged["is_new_enrollment"],
        np.nan,
    )
    merged["lead_to_decision_rate"] = np.where(
        merged["is_lead"] > 0,
        merged["is_decision_14_days"] / merged["is_lead"],
        np.nan,
    )
    merged["lead_to_enroll_rate"] = np.where(
        merged["is_lead"] > 0,
        merged["is_new_enrollment"] / merged["is_lead"],
        np.nan,
    )

    return merged


def compute_scorecard(
    current_month: pd.Series,
    prior_year: pd.Series,
    three_month_trend: pd.Series,
    forecast: pd.Series,
) -> pd.DataFrame:
    \"\"\"Build a channel scorecard comparing actuals to benchmarks.

    Each input is a Series with index = metric names
    (spend, leads, decisions_14d, cpl, cpd, cpne, lead_to_enroll_rate).
    \"\"\"
    scorecard = pd.DataFrame({
        "prior_year": prior_year,
        "three_month_trend": three_month_trend,
        "forecast": forecast,
        "actual": current_month,
    })

    scorecard["delta_to_py"] = np.where(
        scorecard["prior_year"] != 0,
        (scorecard["actual"] - scorecard["prior_year"]) / scorecard["prior_year"],
        np.nan,
    )
    scorecard["delta_to_trend"] = np.where(
        scorecard["three_month_trend"] != 0,
        (scorecard["actual"] - scorecard["three_month_trend"]) / scorecard["three_month_trend"],
        np.nan,
    )
    scorecard["delta_to_forecast"] = np.where(
        scorecard["forecast"] != 0,
        (scorecard["actual"] - scorecard["forecast"]) / scorecard["forecast"],
        np.nan,
    )

    return scorecard
"""

# ─── src/transforms/comparisons.py ───
files['src/transforms/comparisons.py'] = """\"\"\"Period-over-period comparison utilities.\"\"\"

import pandas as pd
import numpy as np


def compute_deltas(
    actuals: pd.DataFrame,
    forecast: pd.DataFrame,
    on: list[str] = None,
) -> pd.DataFrame:
    \"\"\"Compute actual vs forecast deltas.

    Args:
        actuals: DataFrame with actual metrics
        forecast: DataFrame with forecasted metrics
        on: Join columns (default: [\"channel\"])
    \"\"\"
    if on is None:
        on = ["channel"]

    merged = pd.merge(actuals, forecast, on=on, suffixes=("_actual", "_forecast"))

    numeric_cols = [c for c in actuals.columns if c not in on and actuals[c].dtype in ["float64", "int64"]]
    for col in numeric_cols:
        actual_col = f"{col}_actual"
        forecast_col = f"{col}_forecast"
        if actual_col in merged.columns and forecast_col in merged.columns:
            merged[f"{col}_delta"] = merged[actual_col] - merged[forecast_col]
            merged[f"{col}_delta_pct"] = np.where(
                merged[forecast_col] != 0,
                merged[f"{col}_delta"] / merged[forecast_col],
                np.nan,
            )

    return merged


def compute_ytd(metrics: pd.DataFrame, fiscal_year_start: int = 7) -> pd.DataFrame:
    \"\"\"Aggregate metrics year-to-date based on fiscal year start month.\"\"\"
    metrics = metrics.copy()
    metrics["month_num"] = pd.to_datetime(metrics["month_id"], format="%Y%m").dt.month
    metrics["fiscal_year"] = pd.to_datetime(metrics["month_id"], format="%Y%m").dt.year
    mask = metrics["month_num"] >= fiscal_year_start
    metrics.loc[~mask, "fiscal_year"] = metrics.loc[~mask, "fiscal_year"] - 1

    sum_cols = ["spend", "is_lead", "is_application", "is_decision_14_days", "is_new_enrollment"]
    available = [c for c in sum_cols if c in metrics.columns]

    ytd = metrics.groupby(["fiscal_year", "channel"])[available].sum().reset_index()
    return ytd
"""

# ─── src/charts/__init__.py ───
files['src/charts/__init__.py'] = '''"""Chart generation for the monthly report."""

from .cpd_trending import create_cpd_trending_chart
from .actuals_vs_forecast import create_actuals_vs_forecast_chart
from .lead_summary import create_lead_summary_chart

__all__ = [
    "create_cpd_trending_chart",
    "create_actuals_vs_forecast_chart",
    "create_lead_summary_chart",
]
'''

# ─── src/charts/cpd_trending.py ───
files['src/charts/cpd_trending.py'] = '''"""Cost Per Decision trending chart for each channel."""

import logging
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import pandas as pd

logger = logging.getLogger(__name__)

UAGC_BLUE = "#1C2B4A"
UAGC_RED = "#C8102E"
TARGET_GREEN = "#2ECC71"


def create_cpd_trending_chart(
    metrics: pd.DataFrame,
    channel: str,
    target_cpd: float | None = None,
    output_path: Path | str = None,
) -> Path:
    """Generate a CPD trending line chart for a given channel.

    Args:
        metrics: DataFrame with columns [month_id, channel, cpd]
        channel: Channel name to filter
        target_cpd: Optional target line value
        output_path: Where to save the PNG
    """
    channel_data = metrics[metrics["channel"] == channel].sort_values("month_id")

    if channel_data.empty:
        logger.warning("No data for channel: %s", channel)
        return None

    fig, ax = plt.subplots(figsize=(10, 4.5))

    months = channel_data["month_id"].astype(str)
    month_labels = pd.to_datetime(months, format="%Y%m").strftime("%b %y")
    cpd_values = channel_data["cpd"]

    ax.plot(month_labels, cpd_values, color=UAGC_BLUE, linewidth=2.5, marker="o", markersize=5)

    if target_cpd is not None:
        ax.axhline(y=target_cpd, color=TARGET_GREEN, linestyle="--", linewidth=1.5, label="Target")
        ax.legend(loc="upper right", fontsize=9)

    ax.set_title(f"{channel} — Cost Per Decision Trending", fontsize=13, fontweight="bold", color=UAGC_BLUE)
    ax.yaxis.set_major_formatter(mticker.StrMethodFormatter("${x:,.0f}"))
    ax.set_xlabel("")
    ax.set_ylabel("Cost Per Decision", fontsize=10)
    ax.tick_params(axis="x", rotation=45, labelsize=9)
    ax.tick_params(axis="y", labelsize=9)
    ax.grid(axis="y", alpha=0.3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    plt.tight_layout()

    if output_path is None:
        output_path = Path(f"output/cpd_trending_{channel.lower().replace(' ', '_')}.png")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, dpi=150, bbox_inches="tight")
    plt.close(fig)

    logger.info("Saved CPD chart: %s", output_path)
    return output_path
'''

# ─── src/charts/actuals_vs_forecast.py ───
files['src/charts/actuals_vs_forecast.py'] = '''"""Actuals vs Forecast bar chart (Investment & Leads)."""

import logging
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd

logger = logging.getLogger(__name__)

UAGC_BLUE = "#1C2B4A"
UAGC_RED = "#C8102E"
LIGHT_BLUE = "#4A90D9"


def create_actuals_vs_forecast_chart(
    data: pd.DataFrame,
    report_month_label: str,
    output_path: Path | str = None,
) -> Path:
    """Generate grouped bar chart showing spend & leads: forecast vs actual by channel.

    Args:
        data: DataFrame with columns [channel, spend_forecast, spend_actual, leads_forecast, leads_actual]
        report_month_label: e.g. "April 2026"
        output_path: Where to save the PNG
    """
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

    channels = data["channel"].tolist()
    x = np.arange(len(channels))
    width = 0.35

    # Spend chart
    ax1.bar(x - width / 2, data["spend_forecast"], width, label="Forecast", color=LIGHT_BLUE)
    ax1.bar(x + width / 2, data["spend_actual"], width, label="Actual", color=UAGC_BLUE)
    ax1.set_title(f"{report_month_label} Advertising Investment", fontsize=12, fontweight="bold")
    ax1.set_xticks(x)
    ax1.set_xticklabels(channels, rotation=30, ha="right", fontsize=9)
    ax1.yaxis.set_major_formatter(mticker.StrMethodFormatter("${x:,.0f}"))
    ax1.legend(fontsize=9)
    ax1.grid(axis="y", alpha=0.3)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)

    # Leads chart
    ax2.bar(x - width / 2, data["leads_forecast"], width, label="Forecast", color=LIGHT_BLUE)
    ax2.bar(x + width / 2, data["leads_actual"], width, label="Actual", color=UAGC_BLUE)
    ax2.set_title(f"{report_month_label} Paid Leads", fontsize=12, fontweight="bold")
    ax2.set_xticks(x)
    ax2.set_xticklabels(channels, rotation=30, ha="right", fontsize=9)
    ax2.yaxis.set_major_formatter(mticker.StrMethodFormatter("{x:,.0f}"))
    ax2.legend(fontsize=9)
    ax2.grid(axis="y", alpha=0.3)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)

    plt.tight_layout()

    if output_path is None:
        output_path = Path("output/actuals_vs_forecast.png")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, dpi=150, bbox_inches="tight")
    plt.close(fig)

    logger.info("Saved actuals vs forecast chart: %s", output_path)
    return output_path
'''

# ─── src/charts/lead_summary.py ───
files['src/charts/lead_summary.py'] = '''"""Lead to Decision summary table/chart for the report."""

import logging
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd

logger = logging.getLogger(__name__)

UAGC_BLUE = "#1C2B4A"
UAGC_RED = "#C8102E"


def create_lead_summary_chart(
    data: pd.DataFrame,
    report_month_label: str,
    output_path: Path | str = None,
) -> Path:
    """Generate the lead-to-decision summary table as an image.

    Args:
        data: DataFrame with columns [channel, forecast_spend, actual_spend, spend_delta,
              forecast_leads, actual_leads, lead_delta, target_cpl, actual_cpl,
              target_lead_to_enroll, expected_enrolls]
        report_month_label: e.g. "April 2026"
        output_path: Where to save the PNG
    """
    fig, ax = plt.subplots(figsize=(12, 3 + len(data) * 0.4))
    ax.axis("off")

    col_labels = [
        "Channel", "Forecast $", "Actual $", "Delta %",
        "Fcst Leads", "Actual Leads", "Delta %",
        "Target CPL", "Actual CPL",
    ]
    table_data = []
    for _, row in data.iterrows():
        table_data.append([
            row["channel"],
            f"${row.get('spend_forecast', 0):,.0f}",
            f"${row.get('spend_actual', 0):,.0f}",
            f"{row.get('spend_delta_pct', 0):.1%}",
            f"{row.get('leads_forecast', 0):,.0f}",
            f"{row.get('leads_actual', 0):,.0f}",
            f"{row.get('leads_delta_pct', 0):.1%}",
            f"${row.get('target_cpl', 0):,.0f}",
            f"${row.get('actual_cpl', 0):,.0f}",
        ])

    table = ax.table(
        cellText=table_data,
        colLabels=col_labels,
        loc="center",
        cellLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1.2, 1.4)

    # Style header
    for j in range(len(col_labels)):
        table[0, j].set_facecolor(UAGC_BLUE)
        table[0, j].set_text_props(color="white", fontweight="bold")

    ax.set_title(
        f"{report_month_label} — Lead Summary",
        fontsize=13, fontweight="bold", color=UAGC_BLUE, pad=20,
    )

    plt.tight_layout()

    if output_path is None:
        output_path = Path("output/lead_summary.png")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, dpi=150, bbox_inches="tight")
    plt.close(fig)

    logger.info("Saved lead summary chart: %s", output_path)
    return output_path
'''

# ─── src/pptx_builder/__init__.py ───
files['src/pptx_builder/__init__.py'] = '''"""PowerPoint report builder."""

from .builder import ReportBuilder

__all__ = ["ReportBuilder"]
'''

# ─── src/pptx_builder/builder.py ───
files['src/pptx_builder/builder.py'] = '''"""Builds the monthly Performance Marketing PowerPoint from template + data."""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class ReportBuilder:
    """Assembles the monthly report PowerPoint.

    Uses a template PPTX as the base (preserving branding, layouts, master slides)
    and populates it with generated charts, tables, and placeholder text.
    """

    def __init__(self, template_path: Path | str):
        self.template_path = Path(template_path)
        self.prs = Presentation(str(self.template_path))
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

    def set_title_slide(self, month_label: str):
        """Update the title slide with the report month."""
        slide = self.prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if any(m in run.text for m in [
                            "January", "February", "March", "April", "May", "June",
                            "July", "August", "September", "October", "November", "December"
                        ]):
                            run.text = month_label

    def replace_image_on_slide(self, slide_index: int, image_path: Path, image_name: str = None):
        """Replace the first image (or named image) on a slide with a new one."""
        slide = self.prs.slides[slide_index]
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if image_name and shape.name != image_name:
                    continue
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(str(image_path), left, top, width, height)
                logger.info("Replaced image on slide %d", slide_index + 1)
                return True

        logger.warning("No image found on slide %d to replace", slide_index + 1)
        return False

    def update_text_placeholder(self, slide_index: int, placeholder_name: str, new_text: str):
        """Update text in a named placeholder or text box."""
        slide = self.prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == placeholder_name:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = new_text
                return True
        return False

    def update_table_on_slide(self, slide_index: int, data: list[list[str]], table_name: str = None):
        """Update table contents on a slide."""
        slide = self.prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_table:
                if table_name and shape.name != table_name:
                    continue
                table = shape.table
                for row_idx, row_data in enumerate(data):
                    if row_idx >= len(table.rows):
                        break
                    for col_idx, cell_value in enumerate(row_data):
                        if col_idx >= len(table.columns):
                            break
                        table.cell(row_idx, col_idx).text = str(cell_value)
                return True
        return False

    def save(self, output_path: Path | str):
        """Save the assembled presentation."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("Saved report: %s", output_path)
'''

# ─── src/cli.py ───
files['src/cli.py'] = '''"""CLI entry point for generating the monthly report."""

import logging
import sys
from datetime import datetime
from pathlib import Path

import click
import yaml

from .data import BigQuerySource, SQLServerSource
from .transforms.metrics import compute_channel_metrics, compute_scorecard
from .transforms.comparisons import compute_deltas
from .charts import (
    create_cpd_trending_chart,
    create_actuals_vs_forecast_chart,
    create_lead_summary_chart,
)
from .pptx_builder import ReportBuilder

logger = logging.getLogger(__name__)


def load_config(config_path: Path) -> dict:
    """Load YAML configuration."""
    with open(config_path) as f:
        return yaml.safe_load(f)


@click.command()
@click.option("--month", required=True, help="Report month as YYYY-MM (e.g. 2026-04)")
@click.option("--config", default="config/report_config.yaml", help="Path to config YAML")
@click.option("--output", default=None, help="Output PPTX path (default: output/<Month Year>.pptx)")
@click.option("--skip-data", is_flag=True, help="Skip data pull (use cached data)")
@click.option("--verbose", "-v", is_flag=True, help="Enable debug logging")
def main(month: str, config: str, output: str, skip_data: bool, verbose: bool):
    """Generate the UAGC Performance Marketing Monthly Report."""
    log_level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=log_level, format="%(levelname)s | %(name)s | %(message)s")

    config_path = Path(config)
    if not config_path.exists():
        logger.error("Config file not found: %s", config_path)
        sys.exit(1)

    cfg = load_config(config_path)
    report_date = datetime.strptime(month, "%Y-%m")
    month_label = report_date.strftime("%B %Y")
    month_id = report_date.strftime("%Y%m")

    logger.info("Generating report for: %s", month_label)

    # ── Step 1: Pull data ──
    if not skip_data:
        logger.info("Pulling spend data from BigQuery...")
        bq = BigQuerySource(project_id=cfg["bigquery"]["project_id"])
        spend_df = bq.get_channel_spend(month)

        logger.info("Pulling funnel data from SQL Server...")
        sql = SQLServerSource(
            server=cfg["sql_server"]["server"],
            database=cfg["sql_server"]["database"],
            driver=cfg["sql_server"]["driver"],
            trusted_connection=cfg["sql_server"]["trusted_connection"],
        )
        start_month = f"{report_date.year - 1}-{report_date.month:02d}"
        funnel_df = sql.get_funnel_by_channel_month(start_month)
        forecast_df = sql.get_forecast(month_id)
        fraud_df = sql.get_fraud_data(start_month)
    else:
        logger.info("Skipping data pull (--skip-data)")
        # TODO: Load from cached parquet files
        raise NotImplementedError("Cached data loading not yet implemented")

    # ── Step 2: Compute metrics ──
    logger.info("Computing metrics...")
    channel_metrics = compute_channel_metrics(spend_df, funnel_df)

    # ── Step 3: Generate charts ──
    output_dir = Path(cfg.get("output_dir", "output"))
    output_dir.mkdir(parents=True, exist_ok=True)

    logger.info("Generating charts...")
    chart_paths = {}

    for channel in cfg["channels"]:
        display_name = cfg["channel_mapping"].get(channel, channel)
        path = create_cpd_trending_chart(
            channel_metrics,
            channel=display_name,
            output_path=output_dir / f"cpd_trending_{channel}.png",
        )
        if path:
            chart_paths[f"cpd_{channel}"] = path

    avf_chart = create_actuals_vs_forecast_chart(
        compute_deltas(
            channel_metrics[channel_metrics["month_id"] == month_id],
            forecast_df,
        ),
        report_month_label=month_label,
        output_path=output_dir / "actuals_vs_forecast.png",
    )
    chart_paths["actuals_vs_forecast"] = avf_chart

    # ── Step 4: Build PowerPoint ──
    template_path = Path(cfg.get("pptx_template", "templates/template.pptx"))
    if not template_path.exists():
        logger.error("Template not found: %s", template_path)
        sys.exit(1)

    logger.info("Building PowerPoint...")
    builder = ReportBuilder(template_path)
    builder.set_title_slide(month_label)

    # Replace chart images on appropriate slides
    # (Slide indices will be configured per-template)
    # builder.replace_image_on_slide(3, chart_paths["actuals_vs_forecast"])
    # for i, channel in enumerate(cfg["channels"][:4]):
    #     builder.replace_image_on_slide(13 + i*2, chart_paths[f"cpd_{channel}"])

    if output is None:
        output = output_dir / f"{month_label} Monthly Performance Marketing Summary.pptx"

    builder.save(output)
    logger.info("Done! Report saved to: %s", output)


if __name__ == "__main__":
    main()
'''

# ─── Write all files ───
for rel_path, content in files.items():
    full_path = BASE / rel_path
    full_path.parent.mkdir(parents=True, exist_ok=True)
    full_path.write_text(content, encoding='utf-8')
    print(f"  Created: {rel_path}")

print(f"\nDone! {len(files)} files created.")
