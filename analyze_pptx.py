from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r'c:\Users\kseaman\OneDrive - UAGC\General - Performance Marketing\Performance Marketing\Reporting Templates\MonthlyReporting\2026 - 04\April 2026 Monthly Performance Marketing Summary.pptx')

print(f'Total slides: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width} ({Emu(prs.slide_width).inches:.1f} in)')
print(f'Slide height: {prs.slide_height} ({Emu(prs.slide_height).inches:.1f} in)')
print()

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
    print(f'=== Slide {i} (Layout: {layout_name}) ===')
    for shape in slide.shapes:
        stype = str(shape.shape_type).split('(')[0].strip() if '(' in str(shape.shape_type) else str(shape.shape_type)
        print(f'  [{stype}] name="{shape.name}"')
        if shape.has_text_frame:
            text = shape.text_frame.text[:200].replace('\n', ' | ')
            if text.strip():
                print(f'    text: "{text}"')
        if shape.has_table:
            tbl = shape.table
            print(f'    TABLE: {len(tbl.rows)} rows x {len(tbl.columns)} cols')
            for r_idx, row in enumerate(tbl.rows):
                if r_idx < 3:
                    cells = [cell.text[:30] for cell in row.cells]
                    print(f'      row {r_idx}: {cells}')
                elif r_idx == 3:
                    print(f'      ... ({len(tbl.rows) - 3} more rows)')
                    break
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f'    [IMAGE: {shape.image.content_type}]')
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart = shape.chart
            print(f'    [CHART: {chart.chart_type}]')
    print()
